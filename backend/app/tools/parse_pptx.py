"""Parse .pptx -> text via python-pptx."""
from __future__ import annotations

import os
from pptx import Presentation


def parse_pptx(path: str) -> dict:
  """Extract all slide text (titles + bodies + notes)."""
  if not os.path.isfile(path):
    raise FileNotFoundError(path)
  pres = Presentation(path)
  slides = []
  for i, slide in enumerate(pres.slides, 1):
    parts = []
    if slide.shapes.title and slide.shapes.title.text.strip():
      parts.append(slide.shapes.title.text.strip())
    for shape in slide.shapes:
      if shape == slide.shapes.title:
        continue
      if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
          t = "".join(run.text for run in para.runs).strip()
          if t:
            parts.append(t)
      elif shape.has_table:
        for row in shape.table.rows:
          cells = [c.text.strip() for c in row.cells if c.text.strip()]
          if cells:
            parts.append(" | ".join(cells))
    if slide.has_notes_slide:
      notes = slide.notes_slide.notes_text_frame.text.strip()
      if notes:
        parts.append(f"[notes] {notes}")
    if parts:
      slides.append(f"[Slide {i}]\n" + "\n".join(parts))
  content = "\n\n".join(slides).strip()
  title = os.path.splitext(os.path.basename(path))[0] or "Untitled Slides"
  if not content:
    content = "(empty slides)"
  return {"title": title[:200], "content": content}
